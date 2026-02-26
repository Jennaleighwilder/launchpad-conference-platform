#!/usr/bin/env python3
"""Build Launchpad Investor Deck v3 — 12 slides, capital projections, no quantum focus."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

BG = RGBColor(10, 10, 10)

def set_dark_bg(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG

def add_slide(prs, title, body_lines=None, subtitle=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_dark_bg(slide)
    # Title
    t = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    tf = t.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(79, 255, 223)
    # Body
    if body_lines:
        b = slide.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(9), Inches(5))
        bf = b.text_frame
        for i, line in enumerate(body_lines):
            if i == 0:
                p = bf.paragraphs[0]
            else:
                p = bf.add_paragraph()
            p.text = line
            p.font.size = Pt(18)
            p.font.color.rgb = RGBColor(245, 245, 245)
            p.space_after = Pt(8)
    if subtitle:
        s = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(9), Inches(0.8))
        sf = s.text_frame
        sf.paragraphs[0].text = subtitle
        sf.paragraphs[0].font.size = Pt(14)
        sf.paragraphs[0].font.color.rgb = RGBColor(136, 136, 136)
    return slide

def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title
    s1 = prs.slides.add_slide(prs.slide_layouts[6])
    set_dark_bg(s1)
    t = s1.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1.5))
    tf = t.text_frame
    tf.paragraphs[0].text = "Launchpad"
    tf.paragraphs[0].font.size = Pt(54)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = RGBColor(79, 255, 223)
    st = s1.shapes.add_textbox(Inches(1), Inches(4.2), Inches(8), Inches(0.8))
    st.text_frame.paragraphs[0].text = "AI-Powered Event Operating System — 60-Second Conference Generation"
    st.text_frame.paragraphs[0].font.size = Pt(20)
    st.text_frame.paragraphs[0].font.color.rgb = RGBColor(136, 136, 136)

    # Slide 2: Problem
    add_slide(prs, "Problem",
        ["• 6+ months to plan a conference — venue hunting, speaker outreach, schedule drafting",
         "• $50K–150K coordination costs for mid-size events",
         "• 95% of event professionals expect increased AI adoption (Bizzabo 2026)",
         "• Bottleneck: coordination, not ideas"])

    # Slide 3: Solution
    add_slide(prs, "Solution",
        ["• 60 seconds: topic, city, date → full conference page",
         "• 5-agent parallel generation: Speaker, Venue, Schedule, Pricing, Branding",
         "• 6 AI bots promote events: Social Blitz, Community Infiltrator, Email Outreach, Partner Network, SEO & Content, Retargeting",
         "• 34 pages, 11K+ lines, 11 APIs — production-ready"])

    # Slide 4: Product
    add_slide(prs, "Product",
        ["• Event creation wizard: topic, city, date, capacity, budget, vibe",
         "• 6 demo events: AI Summit, Demo Conference, CyberNova, Startup Zaken, etc.",
         "• Ticket checkout (Stripe), sponsor inquiry, affiliate program",
         "• Accessibility: Dyslexia Mode, ADHD Focus, High Contrast, Reduced Motion"])

    # Slide 5: Market Opportunity
    add_slide(prs, "Market Opportunity",
        ["• $2.1T global events industry by 2032 (Allied Market Research)",
         "• 6.4% CAGR growth",
         "• 54% of attendees prefer more in-person events post-pandemic (Eventgroove)",
         "• 95% of event professionals expect AI adoption (Bizzabo 2026)"])

    # Slide 6: Business Model
    add_slide(prs, "Business Model",
        ["• Free: 1 event, AI generation, shareable page",
         "• Pro ($29/mo): Unlimited events, custom branding, Supabase persistence",
         "• Agency: $29/mo, unlimited client events",
         "• Enterprise: Custom pricing, SSO, SLA, dedicated support"])

    # Slide 7: Promotion Engine
    add_slide(prs, "Promotion Engine — 6 AI Street Bots",
        ["• Social Blitz — Twitter, LinkedIn, Instagram posts",
         "• Community Infiltrator — Reddit, Discord, Slack",
         "• Email Outreach — targeted organizer lists",
         "• Partner Network — co-marketing, affiliates",
         "• SEO & Content — blog, landing pages, backlinks",
         "• Retargeting — ads, remarketing campaigns"])

    # Slide 8: Capital Requirements
    add_slide(prs, "Capital Requirements",
        ["• $750K seed round",
         "• 18-month runway",
         "• Allocation: Engineering 40% | Growth 25% | AI Infra 20% | Ops 15%",
         "• Use of funds: team, AI infra, go-to-market, operations"])

    # Slide 9: MRR Projections
    add_slide(prs, "MRR Projections",
        ["• Month 6: $5K MRR",
         "• Month 12: $50K MRR",
         "• Month 18: $150K MRR",
         "• Month 24: $250K MRR",
         "• Assumptions: 12% free-to-paid conversion, $29 Pro avg"])

    # Slide 10: Milestones
    add_slide(prs, "Milestones",
        ["• Q1 2026: 100 paying customers",
         "• Q2 2026: Enterprise pilot (2–3 customers)",
         "• Q3 2026: $50K MRR, Series A prep",
         "• Q4 2026: API launch, integrations"])

    # Slide 11: Team / Ask
    add_slide(prs, "The Ask",
        ["• $750K seed to reach 18-month runway",
         "• Build: 2 engineers, 1 growth, AI infra",
         "• Target: $50K MRR by Month 12",
         "• Contact: [your contact]"])

    # Slide 12: Thank you
    s12 = prs.slides.add_slide(prs.slide_layouts[6])
    set_dark_bg(s12)
    t = s12.shapes.add_textbox(Inches(2), Inches(2.5), Inches(6), Inches(1))
    t.text_frame.paragraphs[0].text = "Thank you"
    t.text_frame.paragraphs[0].font.size = Pt(48)
    t.text_frame.paragraphs[0].font.color.rgb = RGBColor(79, 255, 223)
    st = s12.shapes.add_textbox(Inches(2), Inches(4), Inches(6), Inches(1))
    st.text_frame.paragraphs[0].text = "Launchpad — AI Event Generation"
    st.text_frame.paragraphs[0].font.size = Pt(20)
    st.text_frame.paragraphs[0].font.color.rgb = RGBColor(136, 136, 136)

    out = "Launchpad-Investor-Deck-v3.pptx"
    import os
    out = os.path.join(os.path.dirname(__file__), "..", out)
    prs.save(out)
    print(f"Saved {out}")

if __name__ == "__main__":
    main()
